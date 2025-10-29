from typing import List, Optional
import os
import uuid
from datetime import datetime
from fastapi import APIRouter, HTTPException, UploadFile, File, Depends, Form
from fastapi.responses import Response
from pydantic import BaseModel
try:
    import PyPDF2
except ImportError:
    try:
        import pypdf as PyPDF2
    except ImportError:
        PyPDF2 = None
try:
    import docx
except ImportError:
    try:
        from docx import Document as docx_Document
        docx = type('docx', (), {'Document': docx_Document})
    except ImportError:
        docx = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
import io

# Import services with error handling to prevent startup blocking
try:
    from app.services.storage import StorageService
    print("✓ StorageService imported successfully")
except Exception as e:
    print(f"⚠ StorageService import failed: {e}")
    StorageService = None

try:
    from app.services.summarizer import SummarizerService  
    print("✓ SummarizerService imported successfully")
except Exception as e:
    print(f"⚠ SummarizerService import failed: {e}")
    SummarizerService = None

try:
    from app.services.tts import MultilingualTTSService
    print("✓ MultilingualTTSService imported successfully") 
except Exception as e:
    print(f"⚠ MultilingualTTSService import failed: {e}")
    MultilingualTTSService = None

try:
    from app.utils.security import get_current_user
    print("✓ Security utils imported successfully")
except Exception as e:
    print(f"⚠ Security utils import failed: {e}")
    get_current_user = None

router = APIRouter()

# Lazy-loaded services to speed up startup
_storage = None
_summarizer_service = None
_tts_service = None

def get_storage():
    global _storage
    if _storage is None and StorageService is not None:
        try:
            _storage = StorageService()
        except Exception as e:
            print(f"Failed to initialize StorageService: {e}")
            return None
    return _storage

def get_summarizer_service():
    global _summarizer_service
    if _summarizer_service is None and SummarizerService is not None:
        try:
            _summarizer_service = SummarizerService()
        except Exception as e:
            print(f"Failed to initialize SummarizerService: {e}")
            return None
    return _summarizer_service

def get_tts_service():
    global _tts_service
    if _tts_service is None and MultilingualTTSService is not None:
        try:
            _tts_service = MultilingualTTSService()
        except Exception as e:
            print(f"Failed to initialize MultilingualTTSService: {e}")
            return None
    return _tts_service

@router.get("/health")
async def documents_health():
    """Health check for documents router"""
    return {"status": "ok", "message": "Documents router is working"}

# Old translate endpoint removed - using the improved version below with Google Translate and MyMemory support

@router.get("/test-length/{length}")
async def test_length_calculation(length: int):
    """Test endpoint to verify length calculation"""
    sample_text = """This is a comprehensive business analysis document containing detailed information about strategic planning, operational frameworks, and performance measurement methodologies. The document provides extensive coverage of various business aspects including market analysis, competitive positioning, financial performance indicators, and organizational development strategies. It includes detailed explanations of implementation methodologies, risk management frameworks, and stakeholder engagement processes. The content covers multiple sections including executive summary, methodology, detailed analysis, findings, recommendations, and supporting data. Each section provides comprehensive insights into the subject matter with supporting evidence, examples, and case studies. The document demonstrates thorough examination of best practices, industry standards, and regulatory requirements. It includes analysis of market trends, customer behavior patterns, and competitive landscape dynamics. The methodology section outlines research approaches, data collection methods, and analytical frameworks used in the study. Key findings highlight significant trends, patterns, and insights derived from the comprehensive analysis. The recommendations section provides actionable strategies, implementation roadmaps, and success metrics for organizational improvement. Supporting data includes statistical analysis, performance benchmarks, and comparative studies. The document serves as a comprehensive resource for strategic decision-making, operational planning, and performance optimization. It addresses various stakeholder needs including executive leadership, operational managers, and strategic planning teams. The content is organized in a professional format with clear headings, detailed explanations, and supporting documentation. This comprehensive analysis provides valuable insights for business transformation, process improvement, and strategic growth initiatives."""
    
    extractive_summary = generate_summary(sample_text, "extractive", length)
    abstractive_summary = generate_summary(sample_text, "abstractive", length)
    
    return {
        "length_setting": f"{length}%",
        "original_words": len(sample_text.split()),
        "extractive_summary": extractive_summary,
        "extractive_words": len(extractive_summary.split()),
        "abstractive_summary": abstractive_summary,
        "abstractive_words": len(abstractive_summary.split())
    }

class DocumentResponse(BaseModel):
    id: str
    title: str
    original_filename: str
    content: str
    summary: dict
    tags: List[str]
    created_at: datetime
    file_size: int
    file_type: str

class TranslationRequest(BaseModel):
    text: str
    target_language: str

class TranslationResponse(BaseModel):
    translated_text: str

@router.post("/translate", response_model=TranslationResponse)
async def translate_text(request: TranslationRequest):
    """Translate text using LibreTranslate"""
    import requests
    
    try:
        print(f"=== TRANSLATION REQUEST ===")
        print(f"Target language: {request.target_language}")
        print(f"Text length: {len(request.text)}")
        
        # Try multiple translation services for reliability
        
        # First try: Google Translate API (best for Indian languages)
        try:
            print(f"🌐 Trying Google Translate API for: {request.target_language}")
            google_url = f"https://translate.googleapis.com/translate_a/single?client=gtx&sl=en&tl={request.target_language}&dt=t&q={requests.utils.quote(request.text)}"
            response = requests.get(google_url, timeout=15)
            response.raise_for_status()
            
            result = response.json()
            if result and result[0] and result[0][0] and result[0][0][0]:
                translated_text = ''.join([item[0] for item in result[0] if item[0]])
                if translated_text and len(translated_text.strip()) > 0:
                    print(f"✅ Google Translate success. Output length: {len(translated_text)}")
                    return TranslationResponse(translated_text=translated_text)
        except Exception as e:
            print(f"Google Translate failed: {str(e)}")
        
        # Second try: MyMemory API (good for Indian languages)
        try:
            print(f"🔄 Trying MyMemory API for: {request.target_language}")
            mymemory_url = "https://api.mymemory.translated.net/get"
            params = {
                "q": request.text,
                "langpair": f"en|{request.target_language}"
            }
            response = requests.get(mymemory_url, params=params, timeout=15)
            response.raise_for_status()
            
            result = response.json()
            if result.get("responseStatus") == 200:
                translated_text = result.get("responseData", {}).get("translatedText", "")
                if translated_text and len(translated_text.strip()) > 0:
                    print(f"✅ MyMemory API success. Output length: {len(translated_text)}")
                    return TranslationResponse(translated_text=translated_text)
        except Exception as e:
            print(f"MyMemory API failed: {str(e)}")
        
        # Third try: LibreTranslate - Handle long text by splitting if needed
        try:
            # If text is very long, split into chunks
            if len(request.text) > 2000:
                print(f"Text too long ({len(request.text)} chars), splitting into chunks...")
                
                def translate_libretranslate_chunks(text, target_lang):
                    sentences = text.split('. ')
                    chunks = []
                    current_chunk = ""
                    
                    for sentence in sentences:
                        if len(current_chunk + sentence + '. ') <= 1500:  # LibreTranslate can handle larger chunks
                            current_chunk += sentence + '. '
                        else:
                            if current_chunk:
                                chunks.append(current_chunk.strip())
                            current_chunk = sentence + '. '
                    
                    if current_chunk:
                        chunks.append(current_chunk.strip())
                    
                    translated_chunks = []
                    for i, chunk in enumerate(chunks):
                        print(f"LibreTranslate chunk {i+1}/{len(chunks)} (length: {len(chunk)})")
                        
                        libretranslate_url = "https://libretranslate.de/translate"
                        translation_data = {
                            "q": chunk,
                            "source": "en",
                            "target": target_lang,
                            "format": "text"
                        }
                        
                        response = requests.post(libretranslate_url, json=translation_data, timeout=15)
                        response.raise_for_status()
                        
                        result = response.json()
                        chunk_translation = result.get("translatedText", "")
                        if chunk_translation:
                            translated_chunks.append(chunk_translation)
                        else:
                            translated_chunks.append(chunk)
                    
                    return ' '.join(translated_chunks)
                
                translated_text = translate_libretranslate_chunks(request.text, request.target_language)
            else:
                # For shorter text, translate directly
                libretranslate_url = "https://libretranslate.de/translate"
                translation_data = {
                    "q": request.text,
                    "source": "en",
                    "target": request.target_language,
                    "format": "text"
                }
                
                response = requests.post(libretranslate_url, json=translation_data, timeout=15)
                response.raise_for_status()
                
                result = response.json()
                translated_text = result.get("translatedText", "")
            
            if translated_text and len(translated_text.strip()) > 0:
                print(f"LibreTranslate success. Output length: {len(translated_text)}")
                return TranslationResponse(translated_text=translated_text)
                
        except Exception as e:
            print(f"LibreTranslate failed: {str(e)}")
        
        # Final fallback: MyMemory with chunking for very long text
        try:
            print(f"🔄 Trying MyMemory API with chunking for long text...")
            def translate_chunks(text, target_lang, chunk_size=400):
                """Translate text in chunks to handle length limits"""
                sentences = text.split('. ')
                chunks = []
                current_chunk = ""
                
                for sentence in sentences:
                    if len(current_chunk + sentence + '. ') <= chunk_size:
                        current_chunk += sentence + '. '
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = sentence + '. '
                
                if current_chunk:
                    chunks.append(current_chunk.strip())
                
                translated_chunks = []
                for i, chunk in enumerate(chunks):
                    print(f"Translating chunk {i+1}/{len(chunks)} (length: {len(chunk)})")
                    
                    mymemory_url = "https://api.mymemory.translated.net/get"
                    params = {
                        "q": chunk,
                        "langpair": f"en|{target_lang}"
                    }
                    
                    response = requests.get(mymemory_url, params=params, timeout=15)
                    response.raise_for_status()
                    
                    result = response.json()
                    if result.get("responseStatus") == 200:
                        chunk_translation = result.get("responseData", {}).get("translatedText", "")
                        if chunk_translation:
                            translated_chunks.append(chunk_translation)
                        else:
                            translated_chunks.append(chunk)  # Keep original if translation fails
                    else:
                        translated_chunks.append(chunk)  # Keep original if API fails
                
                return ' '.join(translated_chunks)
            
            translated_text = translate_chunks(request.text, request.target_language)
            
            if translated_text and len(translated_text.strip()) > 0:
                print(f"✅ MyMemory chunk translation success. Output length: {len(translated_text)}")
                return TranslationResponse(translated_text=translated_text)
                    
        except Exception as e:
            print(f"MyMemory chunk translation failed: {str(e)}")
        
        # If all services fail, return error
        raise HTTPException(status_code=503, detail="All translation services unavailable")
        
    except requests.exceptions.RequestException as e:
        print(f"LibreTranslate API error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Translation service error: {str(e)}")
    except Exception as e:
        print(f"Translation error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Translation failed: {str(e)}")

# New multilingual endpoints
@router.get("/languages")
async def get_supported_languages():
    """Get list of all supported languages for multilingual summarization"""
    try:
        languages_info = get_summarizer_service().get_supported_languages()
        return {
            "status": "success",
            "data": languages_info
        }
    except Exception as e:
        return {
            "status": "error",
            "message": f"Failed to load languages: {str(e)}",
            "data": {
                "languages": [
                    {"code": "en", "name": "English", "native": "English"},
                    {"code": "es", "name": "Spanish", "native": "Español"},
                    {"code": "fr", "name": "French", "native": "Français"},
                    {"code": "de", "name": "German", "native": "Deutsch"},
                    {"code": "it", "name": "Italian", "native": "Italiano"},
                    {"code": "pt", "name": "Portuguese", "native": "Português"},
                    {"code": "ru", "name": "Russian", "native": "Русский"},
                    {"code": "zh", "name": "Chinese", "native": "中文"},
                    {"code": "ja", "name": "Japanese", "native": "日本語"},
                    {"code": "ko", "name": "Korean", "native": "한국어"},
                    {"code": "ar", "name": "Arabic", "native": "العربية"},
                    {"code": "hi", "name": "Hindi", "native": "हिन्दी"}
                ],
                "total_count": 12
            }
        }

class MultilingualSummaryRequest(BaseModel):
    text: str
    target_language: str = "en"
    max_length: int = 150
    summary_type: str = "extractive"  # extractive, abstractive, or both

class LanguageDetectionRequest(BaseModel):
    text: str

@router.post("/detect-language")
async def detect_text_language(request: LanguageDetectionRequest):
    """Detect the language of input text"""
    try:
        if not request.text or len(request.text.strip()) < 10:
            return {
                "status": "error",
                "message": "Text too short for reliable language detection",
                "data": {
                    "primary_language": "en",
                    "confidence": 0.3,
                    "possible_languages": [("en", 0.3)],
                    "is_supported": True
                }
            }
        
        detection_result = get_summarizer_service().detect_language(request.text)
        return {
            "status": "success",
            "data": detection_result
        }
    except Exception as e:
        return {
            "status": "error",
            "message": f"Language detection failed: {str(e)}",
            "data": {
                "primary_language": "en",
                "confidence": 0.3,
                "possible_languages": [("en", 0.3)],
                "is_supported": True
            }
        }

@router.post("/summarize-multilingual")
async def create_multilingual_summary(request: MultilingualSummaryRequest):
    """Generate multilingual summary with language detection and translation"""
    try:
        if not request.text or len(request.text.strip()) < 50:
            return {
                "status": "error",
                "message": "Text too short for meaningful summarization",
                "data": {}
            }
        
        # Generate multilingual summary
        summary_result = get_summarizer_service().generate_multilingual_summary(
            text=request.text,
            target_language=request.target_language,
            max_length=request.max_length,
            summary_type=request.summary_type
        )
        
        return {
            "status": "success",
            "data": summary_result
        }
    except Exception as e:
        return {
            "status": "error",
            "message": f"Multilingual summarization failed: {str(e)}",
            "data": {
                "extractive_summary": "Summarization service temporarily unavailable.",
                "abstractive_summary": "Summarization service temporarily unavailable.",
                "detected_language": "en",
                "target_language": request.target_language,
                "confidence": 0.0
            }
        }

@router.post("/upload")
async def upload_document(
    file: UploadFile = File(...),
    summary_type: str = Form(default="extractive"),
    summary_length: int = Form(default=10),
    target_language: str = Form(default="en")
):
    """Upload and process a document with timeout handling"""
    import asyncio
    from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeoutError
    
    print(f"=== UPLOAD REQUEST RECEIVED ===")
    print(f"File: {file.filename}")
    print(f"Content Type: {file.content_type}")
    print(f"Summary Type: {summary_type}")
    print(f"Summary Length: {summary_length}")
    
    # Set timeouts for different operations
    FILE_PROCESSING_TIMEOUT = 30  # 30 seconds for file processing
    AI_PROCESSING_TIMEOUT = 45    # 45 seconds for AI operations
    TOTAL_TIMEOUT = 90           # 90 seconds total
    
    try:
        # Validate file type
        allowed_types = ['.pdf', '.docx', '.txt', '.pptx']
        file_ext = os.path.splitext(file.filename)[1].lower()
        print(f"File extension: {file_ext}")
        
        if file_ext not in allowed_types:
            print(f"ERROR: Unsupported file type: {file_ext}")
            return {"error": f"Unsupported file type: {file_ext}. Allowed types: {allowed_types}"}
        
        # Read file content
        file_content = await file.read()
        print(f"File size: {len(file_content)} bytes")
        
        # Extract text based on file type with timeout
        async def extract_text_with_timeout():
            """Extract text with timeout handling"""
            if file_ext == '.pdf':
                try:
                    print(f"Attempting to extract text from PDF: {file.filename}")
                    
                    # Use thread pool for potentially slow PDF extraction
                    with ThreadPoolExecutor(max_workers=1) as executor:
                        future = executor.submit(extract_pdf_text, file_content)
                        try:
                            text_content = await asyncio.wait_for(
                                asyncio.wrap_future(future), 
                                timeout=FILE_PROCESSING_TIMEOUT
                            )
                        except asyncio.TimeoutError:
                            future.cancel()
                            raise Exception("PDF extraction timed out")
                    
                    print(f"PDF extraction successful! Extracted {len(text_content)} characters")
                    
                    # Validate that we got actual content, not just whitespace
                    if not text_content.strip() or len(text_content.strip()) < 50:
                        raise Exception("Extracted text is too short or empty")
                    
                    return text_content
                    
                except Exception as e:
                    print(f"PDF extraction failed: {str(e)}")
                    # Only use fallback if extraction completely fails
                    return f"""
PDF Extraction Error for "{file.filename}":

The PDF file could not be processed automatically. This could be due to:
- Password protection
- Scanned images without OCR text
- Corrupted file format
- Complex formatting
- Processing timeout

Please try:
1. Converting to a text-based PDF
2. Using a different file format (DOCX, TXT)
3. Ensuring the file is not password protected

Filename: {file.filename}
File size: {len(file_content)} bytes
"""
            elif file_ext == '.docx':
                try:
                    print(f"Attempting to extract text from DOCX: {file.filename}")
                    
                    # Use thread pool for DOCX extraction
                    with ThreadPoolExecutor(max_workers=1) as executor:
                        future = executor.submit(extract_docx_text, file_content)
                        try:
                            text_content = await asyncio.wait_for(
                                asyncio.wrap_future(future), 
                                timeout=FILE_PROCESSING_TIMEOUT
                            )
                        except asyncio.TimeoutError:
                            future.cancel()
                            raise Exception("DOCX extraction timed out")
                    
                    print(f"DOCX extraction successful! Extracted {len(text_content)} characters")
                    
                    # Validate that we got actual content
                    if not text_content.strip() or len(text_content.strip()) < 50:
                        raise Exception("Extracted text is too short or empty")
                    
                    return text_content
                    
                except Exception as e:
                    print(f"DOCX extraction failed: {str(e)}")
                    # Only use fallback if extraction completely fails
                    return f"""
DOCX Extraction Error for "{file.filename}":

The Word document could not be processed automatically. This could be due to:
- Corrupted file format
- Complex formatting or embedded objects
- Password protection
- Unsupported DOCX version
- Processing timeout

Please try:
1. Saving as a simpler DOCX format
2. Converting to TXT format
3. Ensuring the file is not password protected

Filename: {file.filename}
File size: {len(file_content)} bytes
"""
            elif file_ext == '.pptx':
                try:
                    print(f"Attempting to extract text from PPTX: {file.filename}")
                    
                    # Use thread pool for PPTX extraction
                    with ThreadPoolExecutor(max_workers=1) as executor:
                        future = executor.submit(extract_pptx_text, file_content)
                        try:
                            text_content = await asyncio.wait_for(
                                asyncio.wrap_future(future), 
                                timeout=FILE_PROCESSING_TIMEOUT
                            )
                        except asyncio.TimeoutError:
                            future.cancel()
                            raise Exception("PPTX extraction timed out")
                    
                    print(f"PPTX extraction successful! Extracted {len(text_content)} characters")
                    
                    # Validate that we got actual content
                    if not text_content.strip() or len(text_content.strip()) < 50:
                        raise Exception("Extracted text is too short or empty")
                    
                    return text_content
                    
                except Exception as e:
                    print(f"PPTX extraction failed: {str(e)}")
                    # Only use fallback if extraction completely fails
                    return f"""
PPTX Extraction Error for "{file.filename}":

The PowerPoint presentation could not be processed automatically. This could be due to:
- Complex slide layouts with embedded objects
- Password protection
- Corrupted file format
- Slides containing mostly images/graphics
- Processing timeout

Please try:
1. Saving as a simpler PPTX format
2. Converting slide content to DOCX or TXT
3. Ensuring the file is not password protected

Filename: {file.filename}
File size: {len(file_content)} bytes
"""
            elif file_ext == '.txt':
                try:
                    return file_content.decode('utf-8')
                except UnicodeDecodeError:
                    try:
                        return file_content.decode('latin-1')
                    except Exception:
                        return file_content.decode('utf-8', errors='replace')
            else:
                raise Exception(f"Unsupported file type: {file_ext}")
        
        # Execute text extraction with timeout
        try:
            text_content = await asyncio.wait_for(
                extract_text_with_timeout(),
                timeout=FILE_PROCESSING_TIMEOUT + 5
            )
        except asyncio.TimeoutError:
            print("Text extraction timed out")
            return {"error": f"File processing timed out for {file.filename}. Please try a smaller file or different format."}
        
        print(f"Extracted text length: {len(text_content)} characters")
        print(f"First 200 chars: {text_content[:200]}...")
        
        # Generate summary using AI summarization service with timeout
        print(f"=== CALLING AI SUMMARIZATION SERVICE ===")
        print(f"Text length: {len(text_content)}, Summary type: {summary_type}, Length: {summary_length}%")
        
        # Calculate max_length based on percentage (more generous for better summaries)
        base_length = max(150, min(500, len(text_content) // 10))  # Base length between 150-500
        max_length = int(base_length * (summary_length / 100.0) * 2)  # Allow 2x for better quality
        max_length = max(100, min(800, max_length))  # Ensure reasonable bounds
        
        print(f"Calculated max_length: {max_length} characters")
        
        # Use timeout wrapper for AI operations
        async def generate_summaries_with_timeout():
            """Generate summaries with timeout handling"""
            try:
                # Get summarizer with timeout
                summarizer = get_summarizer_service()
                if not summarizer:
                    raise Exception("Summarizer service not available")
                
                # Generate extractive summary with timeout
                with ThreadPoolExecutor(max_workers=1) as executor:
                    future = executor.submit(summarizer.generate_summary, text_content, max_length)
                    try:
                        extractive_summary = await asyncio.wait_for(
                            asyncio.wrap_future(future), 
                            timeout=AI_PROCESSING_TIMEOUT
                        )
                    except asyncio.TimeoutError:
                        future.cancel()
                        raise Exception("Extractive summarization timed out")
                
                print(f"=== EXTRACTIVE SUMMARY GENERATED ===")
                print(f"Length: {len(extractive_summary)} chars")
                print(f"First 100 chars: {extractive_summary[:100]}...")
                
                # For abstractive, try with shorter timeout or fallback
                abstractive_summary = extractive_summary  # Default fallback
                try:
                    # Try multilingual summarization with shorter timeout
                    with ThreadPoolExecutor(max_workers=1) as executor:
                        future = executor.submit(
                            summarizer.generate_multilingual_summary,
                            text_content, 
                            target_language="en",
                            max_length=max_length,
                            summary_type="both"
                        )
                        try:
                            multilingual_result = await asyncio.wait_for(
                                asyncio.wrap_future(future), 
                                timeout=20  # Shorter timeout for abstractive
                            )
                            abstractive_summary = multilingual_result.get('abstractive_summary', extractive_summary)
                        except asyncio.TimeoutError:
                            future.cancel()
                            print("Multilingual summarization timed out, using extractive")
                            abstractive_summary = extractive_summary
                except Exception as e:
                    print(f"Multilingual summarization failed, using extractive: {e}")
                    abstractive_summary = extractive_summary
                
                print(f"=== ABSTRACTIVE SUMMARY GENERATED ===")
                print(f"Length: {len(abstractive_summary)} chars")
                
                return extractive_summary, abstractive_summary
                
            except Exception as e:
                print(f"AI summarization failed: {e}")
                # Fallback to simple text truncation
                fallback_summary = generate_fallback_summary(text_content, max_length)
                return fallback_summary, fallback_summary
        
        # Execute with timeout
        try:
            extractive_summary, abstractive_summary = await asyncio.wait_for(
                generate_summaries_with_timeout(),
                timeout=AI_PROCESSING_TIMEOUT + 10
            )
        except asyncio.TimeoutError:
            print("AI processing timed out completely, using fallback")
            fallback_summary = generate_fallback_summary(text_content, max_length)
            extractive_summary = abstractive_summary = fallback_summary
        
        print(f"Generated extractive summary length: {len(extractive_summary)} characters")
        print(f"Generated abstractive summary length: {len(abstractive_summary)} characters")
        
        # Generate tags using AI keyword extraction
        print(f"=== GENERATING TAGS ===")
        try:
            from app.services.keyword_extractor import KeywordExtractor
            keyword_extractor = KeywordExtractor()
            keywords = keyword_extractor.extract_keywords(text_content, num_keywords=8)
            
            # Extract keyword strings and combine with filename-based tags
            ai_tags = [kw['keyword'].title() for kw in keywords[:5]]
            filename_tags = generate_tags(text_content, file.filename)
            
            # Combine and deduplicate
            all_tags = list(set(ai_tags + filename_tags))[:8]
            tags = all_tags if all_tags else ['Document', 'Analysis']
            
            print(f"Generated {len(tags)} tags: {tags}")
        except Exception as e:
            print(f"AI keyword extraction failed, using fallback: {e}")
            tags = generate_tags(text_content, file.filename)
        
        # Create document record
        document_id = str(uuid.uuid4())
        document_data = {
            "id": document_id,
            "title": file.filename,
            "original_filename": file.filename,
            "content": text_content,
            "summary": {
                "extractive": extractive_summary,
                "abstractive": abstractive_summary
            },
            "tags": tags,
            "created_at": datetime.utcnow(),
            "file_size": len(file_content),
            "file_type": file_ext,
            "owner_id": "test_user",
            "file_data": file_content
        }
        
        # Save to MongoDB with timeout
        print(f"Saving document to MongoDB: {document_data['title']}")
        try:
            saved_id = await asyncio.wait_for(
                get_storage().save_document(document_data),
                timeout=10  # 10 second timeout for database save
            )
            print(f"Document saved with ID: {saved_id}")
        except asyncio.TimeoutError:
            print("Database save timed out")
            return {"error": "Database operation timed out. Please try again."}
        except Exception as e:
            print(f"Database save failed: {e}")
            # Continue anyway, return the data without saving
            saved_id = document_data['id']
        
        # Return response
        return {
            "id": saved_id,
            "title": file.filename,
            "original_filename": file.filename,
            "content": text_content[:500] + "..." if len(text_content) > 500 else text_content,
            "summary": {
                "extractive": extractive_summary,
                "abstractive": abstractive_summary
            },
            "tags": tags,
            "created_at": document_data["created_at"].isoformat(),
            "file_size": len(file_content),
            "file_type": file_ext
        }
        
    except Exception as e:
        print(f"ERROR in upload: {str(e)}")
        import traceback
        traceback.print_exc()
        return {"error": str(e)}

@router.get("/{document_id}", response_model=DocumentResponse)
async def get_document(document_id: str):
    """Get document by ID"""
    document = await get_storage().get_document_by_id(document_id)
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    
    return DocumentResponse(**{k: v for k, v in document.items() if k != "file_data"})

@router.get("/{document_id}/download")
async def download_document(document_id: str):
    """Download original document"""
    print(f"Download request for document ID: {document_id}")
    
    document = await get_storage().get_document_by_id(document_id)
    if not document:
        print(f"Document not found: {document_id}")
        raise HTTPException(status_code=404, detail="Document not found")
    
    print(f"Found document: {document.get('title', 'Unknown')}")
    print(f"File data size: {len(document.get('file_data', b''))} bytes")
    
    if not document.get("file_data"):
        print("No file data found in document")
        raise HTTPException(status_code=404, detail="File data not found")
    
    from fastapi.responses import Response
    
    return Response(
        content=document["file_data"],
        media_type="application/octet-stream",
        headers={"Content-Disposition": f"attachment; filename={document['original_filename']}"}
    )

@router.get("/{document_id}/audio")
async def get_document_audio(document_id: str, language: str = "en"):
    """Generate and return audio for document summary"""
    document = await get_storage().get_document_by_id(document_id)
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Generate audio using text-to-speech
    summary_text = document["summary"]["extractive"]
    audio_data = generate_audio(summary_text, language)
    
    from fastapi.responses import Response
    
    return Response(
        content=audio_data,
        media_type="audio/wav",
        headers={"Content-Disposition": f"attachment; filename=summary_audio.wav"}
    )

def clean_extracted_text(text: str) -> str:
    """Improved text cleaning with proper formatting"""
    import re
    
    try:
        print("=== CLEAN_EXTRACTED_TEXT CALLED ===")
        print(f"Input text length: {len(text)}")
        print(f"Input text first 300 chars: {text[:300]}")
        
        # Basic text cleaning without heavy dependencies
        print("Using improved text cleaning with formatting...")
        
        # Step 0: Fix basic formatting issues first
        # Remove excessive whitespace but preserve paragraph breaks
        text = re.sub(r'[ \t]+', ' ', text)  # Multiple spaces/tabs to single space
        text = re.sub(r'\n[ \t]*\n', '\n\n', text)  # Preserve paragraph breaks
        
        # Step 1: Add proper paragraph breaks for numbered sections
        text = re.sub(r'(\d+\.\d+)\s+([A-Z])', r'\n\n\1 \2', text)  # "3.1 What" -> "\n\n3.1 What"
        text = re.sub(r'(\d+\.)\s+([A-Z])', r'\n\n\1 \2', text)  # "1. Cross" -> "\n\n1. Cross"
        
        # Step 2: Add breaks after periods followed by capital letters (new sentences)
        text = re.sub(r'(\.)([A-Z][a-z])', r'\1 \2', text)  # Ensure space after periods
        
        # Step 3: Fix common formatting issues
        text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)  # Add space between camelCase
        
        # Step 4: Add proper line breaks for lists and bullet points
        text = re.sub(r'(\d+\.\s)', r'\n\1', text)  # Number lists
        text = re.sub(r'([.!?])\s*([A-Z][a-z]+:)', r'\1\n\n\2', text)  # Headers after sentences
        
        # Step 1: Fix ligatures and special characters FIRST
        text = text.replace('ﬁ', 'fi')
        text = text.replace('ﬂ', 'fl')
        text = text.replace('ﬀ', 'ff')
        text = text.replace('ﬃ', 'ffi')
        text = text.replace('ﬄ', 'ffl')
        text = text.replace('ﬆ', 'st')
        
        # Step 1.5: Fix broken words with spaces in the middle
        # Pattern: single letter followed by space followed by rest of word (e.g., "Network s" -> "Networks")
        text = re.sub(r'\b(\w+) ([a-z]{1,3})\b', r'\1\2', text)  # Fix "word s" -> "words", "data ta" -> "data"
        
        # Fix common broken words with multiple spaces
        broken_patterns = [
            # Latest user examples
            (r'\basa\b', 'as a'),
            (r'employsa\b', 'employs a'),
            (r'indicates th\b', 'indicates that'),
            (r'Atthe\b', 'At the'),
            (r'atthe\b', 'at the'),
            
            # From previous user examples
            (r'needfor\b', 'need for'),
            (r'libraryand\b', 'library and'),
            (r'S cal ability', 'Scalability'),
            (r'timingsof\b', 'timings of'),
            (r'sharingin\b', 'sharing in'),
            (r'parallel ize', 'parallelize'),
            (r'parallelizethe', 'parallelize the'),
            (r'codeto\b', 'code to'),
            (r'solvethe\b', 'solve the'),
            (r'endof\b', 'end of'),
            (r'willbe\b', 'will be'),
            (r'ableto\b', 'able to'),
            (r'libraryto\b', 'library to'),
            (r'programfor\b', 'program for'),
            (r'math em at ical', 'mathematical'),
            (r'puremath ematics', 'pure mathematics'),
            
            # Previous patterns
            (r'serie s\b', 'series'),
            (r'ex ample', 'example'),
            (r'in clud', 'includ'),
            (r'de scrib', 'describ'),
            (r'moti v ation', 'motivation'),
            (r'oper ation', 'operation'),
            (r'func tion', 'function'),
            (r'deﬁ nition', 'definition'),
            (r'de finition', 'definition'),
            (r'eﬃ cient', 'efficient'),
            (r'ef ficient', 'efficient'),
            (r'diﬀ erent', 'different'),
            (r'dif ferent', 'different'),
            (r'ﬁ rst', 'first'),
            (r' rst\b', 'first'),
            (r'we rst', 'we first'),
            (r'area\b', 'are a'),  # "area" when it should be "are a"
        ]
        
        for pattern, replacement in broken_patterns:
            text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
        
        # Step 2: Preserve paragraph structure while cleaning
        # Keep double line breaks (paragraph boundaries)
        text = re.sub(r'\n\n+', '\n\n', text)  # Normalize paragraph breaks
        # Replace single line breaks with spaces (but preserve paragraphs)
        text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
        # Normalize spaces within lines
        text = re.sub(r'[ \t]+', ' ', text)
        
        # Step 3: Fix obvious issues first
        text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)  # Add space between camelCase
        text = re.sub(r'([A-Za-z])(\d)', r'\1 \2', text)  # Add space between letter and number
        text = re.sub(r'(\d)([A-Za-z])', r'\1 \2', text)  # Add space between number and letter
        text = re.sub(r'\.([A-Z])', r'. \1', text)  # Add space after periods
        text = re.sub(r',([A-Za-z])', r', \1', text)  # Add space after commas
        
        # Step 3: Split text into sentences for processing
        sentences = sent_tokenize(text)
        cleaned_sentences = []
        
        for sentence in sentences:
            # Step 4: Split sentence into words
            words = word_tokenize(sentence)
            cleaned_words = []
            
            for word in words:
                # Remove punctuation for processing but keep track of it
                punct = ''
                clean_word = word
                if word and not word[-1].isalnum():
                    punct = word[-1]
                    clean_word = word[:-1]
                
                # Step 5: Apply WordNinja to EVERY alphabetic word
                # Let WordNinja's AI decide if it should be split
                # Process words >= 3 chars to catch short concatenations like "asa" -> "as a"
                if clean_word.isalpha() and len(clean_word) >= 3:
                    # WordNinja uses statistical NLP models to intelligently split
                    split_words = wordninja.split(clean_word)
                    
                    # Only use the split if it results in multiple words AND makes sense
                    # WordNinja will return single word if it's already correct
                    if len(split_words) > 1:
                        # Check if split makes sense - avoid over-splitting
                        # Allow single-letter words only if they're common: a, I
                        # Otherwise parts should be >= 2 chars
                        common_single_letters = {'a', 'A', 'I'}
                        valid_split = all(
                            len(w) >= 2 or w in common_single_letters 
                            for w in split_words
                        )
                        
                        if valid_split:
                            print(f"WordNinja split: '{clean_word}' -> {split_words}")
                            # Add split words
                            cleaned_words.extend(split_words[:-1])
                            # Add last word with punctuation if any
                            cleaned_words.append(split_words[-1] + punct if punct else split_words[-1])
                        else:
                            print(f"WordNinja split rejected (too short): '{clean_word}' -> {split_words}")
                            # Split created very short words, keep original
                            cleaned_words.append(word)
                    else:
                        # WordNinja says it's already a single word
                        print(f"WordNinja no split needed: '{clean_word}'")
                        cleaned_words.append(word)
                else:
                    # Keep short words or non-alphabetic words as-is
                    cleaned_words.append(word)
            
            # Step 6: Reconstruct sentence with proper spacing
            cleaned_sentence = ' '.join(cleaned_words)
            
            # Step 7: Fix punctuation spacing
            cleaned_sentence = re.sub(r'\s+([,.;:!?])', r'\1', cleaned_sentence)
            cleaned_sentence = re.sub(r'([,.;:!?])([A-Za-z])', r'\1 \2', cleaned_sentence)
            
            cleaned_sentences.append(cleaned_sentence)
        
        # Step 8: Reconstruct full text
        cleaned_text = '. '.join(cleaned_sentences)
        
        # Step 9: Final cleanup
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text)  # Remove extra spaces
        cleaned_text = re.sub(r'\.+', '.', cleaned_text)  # Fix multiple periods
        
        # Step 10: Final formatting improvements for readability
        cleaned_text = improve_text_formatting(cleaned_text)
        
        print(f"AI cleaning completed. Original length: {len(text)}, Cleaned length: {len(cleaned_text)}")
        return cleaned_text.strip()
        
    except ImportError as e:
        print(f"AI libraries not available ({e}), falling back to basic cleaning...")
        # Fallback to basic regex cleaning if AI libraries aren't available
        return basic_text_cleaning(text)
    except Exception as e:
        print(f"AI cleaning failed ({e}), falling back to basic cleaning...")
        return basic_text_cleaning(text)

def improve_text_formatting(text: str) -> str:
    """Improve text formatting for better readability with proper paragraphs"""
    import re
    
    # Step 1: Add proper paragraph breaks for chapters and sections
    text = re.sub(r'(Chapter\s+\d+)', r'\n\n\1', text)  # "Chapter 9" -> "\n\nChapter 9"
    text = re.sub(r'(\d+\.\d+)\s+([A-Z])', r'\n\n\1 \2', text)  # "9.1 The" -> "\n\n9.1 The"
    text = re.sub(r'(\d+\.)\s+([A-Z][a-z]+)', r'\n\n\1 \2', text)  # "1. Cross" -> "\n\n1. Cross"
    
    # Step 2: Add breaks after sentences that end topics (look for capital letter starts)
    text = re.sub(r'([.!?])\s+([A-Z][a-z]+\s+[a-z]+.*?[.!?])\s+([A-Z][A-Z])', r'\1\n\n\2\n\n\3', text)
    
    # Step 3: Add breaks before new topics and important concepts
    text = re.sub(r'([.!?])\s*(Convolutional networks|Neural networks|The name|Examples include|Usually|Research into|In this chapter|We first describe|We then describe)', r'\1\n\n\2', text)
    
    # Step 3.1: Add breaks for specific technical content
    text = re.sub(r'([.!?])\s*(CNNs, are|Convolution is|In many ways)', r'\1\n\n\2', text)
    
    # Step 4: Add breaks for figure references and technical descriptions
    text = re.sub(r'([.!?])\s*(Figure\s+\d+)', r'\1\n\n\2', text)
    text = re.sub(r'([.!?])\s*(Input image:|Output)', r'\1\n\n\2', text)
    
    # Step 5: Add breaks for questions and explanatory sections
    text = re.sub(r'([.!?])\s*(Why|What|How|When|Where)\s+([a-z])', r'\1\n\n\2 \3', text)
    
    # Step 6: Add breaks for contrasting statements and new paragraphs
    text = re.sub(r'([.!?])\s*(However|Nevertheless|In contrast|On the other hand|Meanwhile|Furthermore|Moreover|Additionally)', r'\1\n\n\2', text)
    
    # Step 7: Add breaks for conclusions and transitions
    text = re.sub(r'([.!?])\s*(In conclusion|To conclude|Finally|In summary|Overall|This approach)', r'\1\n\n\2', text)
    
    # Step 8: Fix spacing around periods and ensure proper sentence separation
    text = re.sub(r'([a-z])([A-Z][a-z])', r'\1. \2', text)  # Add missing periods between sentences
    text = re.sub(r'\.([A-Z][a-z])', r'. \1', text)  # Ensure space after periods
    
    # Step 8.1: Add breaks for very long sentences (force paragraph breaks every ~3-4 sentences)
    sentences = text.split('. ')
    formatted_sentences = []
    for i, sentence in enumerate(sentences):
        formatted_sentences.append(sentence)
        # Add paragraph break every 3-4 sentences for better readability
        if (i + 1) % 3 == 0 and i < len(sentences) - 1:
            formatted_sentences.append('\n\n')
    text = '. '.join(formatted_sentences)
    
    # Step 9: Add breaks for technical specifications and lists
    text = re.sub(r'([.!?])\s*(Examples include|The specific|Real convolutional)', r'\1\n\n\2', text)
    
    # Step 10: Clean up excessive line breaks and formatting
    text = re.sub(r'\n{3,}', '\n\n', text)  # Max 2 line breaks
    text = re.sub(r'^\n+', '', text)  # Remove leading line breaks
    text = re.sub(r'\s+', ' ', text)  # Clean up multiple spaces
    text = re.sub(r'\n\s+', '\n', text)  # Remove spaces after line breaks
    
    return text.strip()

def basic_text_cleaning(text: str) -> str:
    """Fallback basic text cleaning with formatting"""
    import re
    
    # Basic regex-based cleaning as fallback
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
    text = re.sub(r'([A-Za-z])(\d)', r'\1 \2', text)
    text = re.sub(r'(\d)([A-Za-z])', r'\1 \2', text)
    text = re.sub(r'\.([A-Z])', r'. \1', text)
    text = re.sub(r',([A-Za-z])', r', \1', text)
    
    # Apply formatting improvements
    text = improve_text_formatting(text)
    
    return text.strip()

def extract_pdf_text(file_content: bytes) -> str:
    """Extract text from PDF using PyMuPDF (superior text extraction)"""
    try:
        # Try PyMuPDF first (much better text extraction)
        import fitz  # PyMuPDF
        
        print("Using PyMuPDF for superior PDF text extraction...")
        
        pdf_file = io.BytesIO(file_content)
        pdf_document = fitz.open(stream=pdf_file, filetype="pdf")
        
        text = ""
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            # Use get_text with "text" mode for best results
            page_text = page.get_text("text")
            if page_text:
                text += page_text + "\n"
        
        pdf_document.close()
        extracted_text = text.strip()
        
        print(f"PyMuPDF extracted {len(extracted_text)} characters from PDF")
        
        # If no text was extracted, fallback to PyPDF2
        if not extracted_text or len(extracted_text.split()) < 5:
            print("PyMuPDF extraction failed, trying PyPDF2...")
            return extract_pdf_text_pypdf2(file_content)
        
        # Clean up the extracted text
        cleaned_text = clean_extracted_text(extracted_text)
        return cleaned_text
        
    except ImportError:
        print("PyMuPDF not available, falling back to PyPDF2...")
        return extract_pdf_text_pypdf2(file_content)
    except Exception as e:
        print(f"PyMuPDF extraction failed ({e}), falling back to PyPDF2...")
        return extract_pdf_text_pypdf2(file_content)

def extract_pdf_text_pypdf2(file_content: bytes) -> str:
    """Fallback PDF extraction using PyPDF2"""
    if PyPDF2 is None:
        raise Exception("PyPDF2 is not installed. Cannot process PDF files.")
    
    try:
        print("Using PyPDF2 for PDF text extraction (fallback)...")
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        
        extracted_text = text.strip()
        
        # If no text was extracted or very little text, provide a meaningful fallback
        if not extracted_text or len(extracted_text.split()) < 5:
            return f"This PDF document appears to contain primarily images, charts, or formatted content that cannot be automatically extracted as text. The document has {len(pdf_reader.pages)} pages and may contain visual elements, diagrams, or special formatting that requires manual review to fully understand the content."
        
        # Clean up the extracted text to fix formatting issues
        cleaned_text = clean_extracted_text(extracted_text)
        return cleaned_text
    except Exception as e:
        raise Exception(f"Error reading PDF: {str(e)}")

def extract_docx_text(file_content: bytes) -> str:
    """Extract text from DOCX"""
    if docx is None:
        raise Exception("python-docx is not installed. Cannot process DOCX files.")
    
    try:
        docx_file = io.BytesIO(file_content)
        doc = docx.Document(docx_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Error reading DOCX: {str(e)}")

def extract_pptx_text(file_content: bytes) -> str:
    """Extract text from PPTX (PowerPoint)"""
    if Presentation is None:
        raise Exception("python-pptx is not installed. Cannot process PPTX files.")
    
    try:
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        text = ""
        
        # Extract text from all slides
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n=== Slide {slide_num} ===\n"
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
                
                # Handle tables separately
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells if cell.text])
                        if row_text:
                            text += row_text + "\n"
            
            text += "\n"
        
        return text.strip()
    except Exception as e:
        raise Exception(f"Error reading PPTX: {str(e)}")

def generate_summary(text: str, summary_type: str, length: int) -> str:
    """Generate summary based on proper summarization standards"""
    try:
        print(f"=== GENERATING {summary_type.upper()} SUMMARY ===")
        print(f"Length setting: {length}%")
        print(f"*** DEBUG: This is the NEW FIXED generate_summary function ***")
        print(f"*** PARAMETERS: text_length={len(text)}, summary_type='{summary_type}', length={length} ***")
    except Exception as e:
        print(f"Error in generate_summary setup: {e}")
        return f"Error generating {summary_type} summary: {str(e)}"
    
    # Clean and prepare text
    text = text.strip()
    if not text:
        return "No content available for summarization."
    
    # Handle very short documents by providing a meaningful summary
    if len(text.split()) < 10:
        return f"This is a brief document containing: {text[:200]}{'...' if len(text) > 200 else ''}"
    
    words = text.split()
    sentences = [s.strip() for s in text.split('.') if s.strip()]
    total_words = len(words)
    total_sentences = len(sentences)
    
    print(f"Original document: {total_words} words, {total_sentences} sentences")
    
    # SIMPLE AND DIRECT length calculation - use percentage more literally
    # The user expects 85% to mean a substantial portion of the document
    
    if summary_type == "extractive":
        # SIMPLE LINEAR CALCULATION: 10% = ~50 words, 100% = ~750 words (3000 chars)
        # This gives a direct relationship between percentage and length
        
        # Calculate target based on linear scale
        # Min: 50 words (10%) = ~200 chars (1-2 sentences)
        # Max: 750 words (100%) = ~3000 chars
        
        min_words = 50  # 1-2 sentences
        max_words = 750  # ~3000 characters
        
        # Linear interpolation: target = min + (percentage/100) * (max - min)
        percentage_decimal = length / 100.0
        target_words = int(min_words + percentage_decimal * (max_words - min_words))
        
        # Ensure we don't exceed document length
        target_words = min(target_words, int(total_words * 0.8))  # Never more than 80% of original
        
        print(f"*** LINEAR CALCULATION: {length}% = {min_words} + {percentage_decimal:.2f} * {max_words - min_words} = {target_words} words ***")
        print(f"Target characters: ~{target_words * 4} chars (aiming for 3000 max)")
    else:
        # Abstractive: Use same linear approach but with sentences
        # Min: 2 sentences (10%), Max: 30 sentences (100%)
        
        min_sentences = 2  # 1-2 sentences
        max_sentences = 30  # ~3000 characters worth
        
        # Linear interpolation for sentences
        percentage_decimal = length / 100.0
        target_sentences = int(min_sentences + percentage_decimal * (max_sentences - min_sentences))
        
        # Ensure we don't exceed document sentences
        target_sentences = min(target_sentences, int(total_sentences * 0.8))  # Never more than 80% of original
        
        print(f"*** LINEAR ABSTRACTIVE: {length}% = {min_sentences} + {percentage_decimal:.2f} * {max_sentences - min_sentences} = {target_sentences} sentences ***")
    
    if summary_type == "extractive":
        # Extractive: Select important sentences and combine
        # Use a smarter approach: take from beginning, middle, and end for better coverage
        
        if target_words >= total_words:
            # If target is larger than original, return most of the text
            summary = ' '.join(words[:int(total_words * 0.95)])
        else:
            # Smart word selection: take from different parts of the document
            # This gives better coverage than just taking from the beginning
            
            # Take 40% from beginning, 30% from middle, 30% from end
            begin_words = int(target_words * 0.4)
            middle_words = int(target_words * 0.3)
            end_words = target_words - begin_words - middle_words
            
            # Beginning section
            beginning = words[:begin_words]
            
            # Middle section
            middle_start = len(words) // 3
            middle_end = middle_start + middle_words
            middle = words[middle_start:middle_end]
            
            # End section
            end_start = max(len(words) - end_words, middle_end + 10)  # Avoid overlap
            ending = words[end_start:]
            
            # Combine sections
            selected_words = beginning + middle + ending
            
            # Ensure we don't exceed target
            if len(selected_words) > target_words:
                selected_words = selected_words[:target_words]
            
            summary = ' '.join(selected_words)
        
        # Debug: Check actual word count
        actual_words = len(summary.split())
        print(f"*** EXTRACTIVE DEBUG: Target={target_words}, Actual={actual_words}, Length={len(summary)} chars ***")
        
        # Only apply character limits for very short summaries
        if length <= 20 and len(summary) > 500:  # Only for 10% summaries
            words_list = summary.split()
            summary = ' '.join(words_list[:50])  # Max 50 words for 10%
            print(f"*** CHARACTER LIMIT: Cut to {len(summary)} chars, {len(summary.split())} words ***")
        
    else:  # abstractive
        # Abstractive: Select key sentences from different parts for better coverage
        if target_sentences >= total_sentences:
            selected_sentences = sentences[:int(total_sentences * 0.95)]  # Use most sentences
        else:
            # Smart sentence selection: take from beginning, middle, and end
            begin_sentences = int(target_sentences * 0.4)
            middle_sentences = int(target_sentences * 0.3)
            end_sentences = target_sentences - begin_sentences - middle_sentences
            
            # Beginning sentences
            beginning = sentences[:begin_sentences]
            
            # Middle sentences
            middle_start = len(sentences) // 3
            middle_end = middle_start + middle_sentences
            middle = sentences[middle_start:middle_end]
            
            # End sentences
            end_start = max(len(sentences) - end_sentences, middle_end + 2)  # Avoid overlap
            ending = sentences[end_start:end_start + end_sentences]
            
            # Combine sections
            selected_sentences = beginning + middle + ending
            
            # Ensure we don't exceed target
            if len(selected_sentences) > target_sentences:
                selected_sentences = selected_sentences[:target_sentences]
        
        summary = '. '.join(selected_sentences)
        
        # Ensure proper punctuation
        if summary and not summary.endswith('.'):
            summary += '.'
        
        # Debug: Check actual sentence count
        final_sentences = len([s for s in summary.split('.') if s.strip()])
        print(f"*** ABSTRACTIVE DEBUG: Target={target_sentences}, Actual={final_sentences}, Length={len(summary)} chars ***")
    
    # Ensure proper punctuation
    if summary and not summary.endswith(('.', '!', '?')):
        summary += '.'
    
    final_words = len(summary.split())
    try:
        print(f"*** FINAL RESULT: Target={target_words if summary_type == 'extractive' else target_sentences}, Actual={final_words} words ***")
        print(f"Generated summary: {final_words} words")
        print(f"Compression ratio: {(final_words/total_words)*100:.1f}%")
        
        # APPLY FORMATTING TO SUMMARY - This is crucial for readability
        print("Applying formatting to summary...")
        formatted_summary = improve_text_formatting(summary)
        print(f"Formatting applied - Original length: {len(summary)}, Formatted length: {len(formatted_summary)}")
        
        return formatted_summary
    except Exception as e:
        print(f"Error in generate_summary final steps: {e}")
        return f"Error completing {summary_type} summary generation: {str(e)}"

def generate_tags(text: str, filename: str) -> List[str]:
    """Generate relevant tags from text and filename"""
    # Simple keyword extraction
    keywords = []
    
    # Extract from filename
    name_parts = filename.lower().replace('.', ' ').replace('_', ' ').replace('-', ' ').split()
    keywords.extend([part.title() for part in name_parts if len(part) > 2])
    
    # Common business/document keywords
    business_keywords = ['business', 'report', 'analysis', 'strategy', 'performance', 'financial', 'quarterly', 'annual']
    text_lower = text.lower()
    
    for keyword in business_keywords:
        if keyword in text_lower:
            keywords.append(keyword.title())
    
    # Remove duplicates and limit to 5 tags
    unique_tags = list(set(keywords))[:5]
    
    # Ensure we have at least some default tags
    if not unique_tags:
        unique_tags = ['Document', 'Analysis']
    
    return unique_tags

class TTSRequest(BaseModel):
    text: str
    language: str = "en"

@router.post("/text-to-speech")
async def text_to_speech(request: TTSRequest):
    """
    Convert text to speech using gTTS (Google Text-to-Speech)
    Supports all 12 core languages: en, es, fr, de, it, pt, ru, zh, ja, ko, ar, hi
    """
    try:
        if not request.text or not request.text.strip():
            raise HTTPException(status_code=400, detail="Text cannot be empty")
        
        # Generate audio using gTTS
        audio_data = tts_service.synthesize(
            text=request.text.strip(),
            lang=request.language
        )
        
        # Return audio as MP3
        return Response(
            content=audio_data,
            media_type="audio/mpeg",
            headers={
                "Content-Disposition": "attachment; filename=speech.mp3",
                "Cache-Control": "no-cache"
            }
        )
        
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(f"TTS Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Text-to-speech failed: {str(e)}")

def generate_fallback_summary(text: str, max_length: int) -> str:
    """Generate a simple fallback summary when AI services fail"""
    try:
        words = text.split()
        sentences = [s.strip() for s in text.split('.') if s.strip()]
        
        # Calculate target length (aim for 20% of original or max_length, whichever is smaller)
        target_words = min(max_length // 4, len(words) // 5, 100)
        target_words = max(target_words, 20)  # Minimum 20 words
        
        if len(words) <= target_words:
            return text.strip()
        
        # Take first portion of text and ensure it ends with a complete sentence
        truncated_words = words[:target_words]
        truncated_text = ' '.join(truncated_words)
        
        # Find the last complete sentence
        last_period = truncated_text.rfind('.')
        if last_period > len(truncated_text) // 2:  # If we found a period in the latter half
            truncated_text = truncated_text[:last_period + 1]
        else:
            # Add ellipsis if no good sentence break found
            truncated_text += "..."
        
        return truncated_text.strip()
        
    except Exception as e:
        print(f"Fallback summary generation failed: {e}")
        # Ultimate fallback - just truncate
        return text[:max_length] + "..." if len(text) > max_length else text

def generate_audio(text: str, language: str = "en") -> bytes:
    """Legacy function - use /text-to-speech endpoint instead"""
    return b"RIFF$\x00\x00\x00WAVEfmt \x10\x00\x00\x00\x01\x00\x01\x00D\xac\x00\x00\x88X\x01\x00\x02\x00\x10\x00data\x00\x00\x00\x00"
